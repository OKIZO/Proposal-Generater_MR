import streamlit as st
import json
import io
import os
from pptx import Presentation
from pptx.util import Inches

# レイアウトを横幅いっぱいに使う設定（2カラムに最適化）
st.set_page_config(page_title="PPTX生成システム", layout="wide")

# ==========================================
# 認証・パスワード管理機能
# ==========================================
CONFIG_FILE = "config.json"
DEFAULT_USER_PWD = "team_creative"
ADMIN_PWD = "okino_creative"

def load_user_pwd():
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                return json.load(f).get("pwd", DEFAULT_USER_PWD)
        except:
            return DEFAULT_USER_PWD
    return DEFAULT_USER_PWD

def save_user_pwd(new_pwd):
    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
        json.dump({"pwd": new_pwd}, f)

if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

# --- ログイン画面の表示（未ログイン時） ---
if not st.session_state.logged_in:
    _, col_center, _ = st.columns([1, 2, 1])
    with col_center:
        st.markdown("<h2 style='text-align:center; margin-top:4rem; margin-bottom:2rem;'>🔐 MedConcept ログイン</h2>", unsafe_allow_html=True)
        pwd_input = st.text_input("チーム用パスワードを入力", type="password")
        if st.button("ログイン", type="primary", use_container_width=True):
            if pwd_input == load_user_pwd():
                st.session_state.logged_in = True
                st.rerun()
            else:
                st.error("パスワードが間違っています。")
        
        st.markdown("<br><br>", unsafe_allow_html=True)
        with st.expander("⚙️ 管理者設定（パスワードの変更）"):
            st.markdown("<small>※管理者のパスワードが必要です</small>", unsafe_allow_html=True)
            admin_input = st.text_input("管理者パスワード", type="password")
            new_pwd_input = st.text_input("新しいチーム用パスワード")
            if st.button("パスワードを更新", use_container_width=True):
                if admin_input == ADMIN_PWD:
                    if new_pwd_input.strip():
                        save_user_pwd(new_pwd_input.strip())
                        st.success(f"チーム用パスワードを「{new_pwd_input.strip()}」に変更しました！")
                    else:
                        st.error("新しいパスワードを入力してください。")
                else:
                    st.error("管理者パスワードが間違っています。")
    st.stop()

# ==========================================
# これより下はログイン成功時のみ実行される
# ==========================================

# --- 補助関数：図形やセル内のテキストをフォント維持で置換 ---
def replace_text_in_shape(item, replacements):
    if not hasattr(item, "text_frame") or item.text_frame is None:
        return
    for paragraph in item.text_frame.paragraphs:
        p_text = "".join(run.text for run in paragraph.runs)
        replaced_any = False
        for old_text, new_text in replacements.items():
            if old_text in p_text:
                p_text = p_text.replace(old_text, str(new_text))
                replaced_any = True
                
        if replaced_any:
            if len(paragraph.runs) > 0:
                paragraph.runs[0].text = p_text
                for i in range(1, len(paragraph.runs)):
                    paragraph.runs[i].text = ""

# --- メイン処理関数 ---
def generate_pptx(json_data, uploaded_images):
    prs = Presentation("template.pptx")

    if "基本情報" in json_data:
        base_info = json_data.get("基本情報", {})
        product_name = base_info.get("製品名", "")
        item_name = base_info.get("アイテム名", "")
        spec = base_info.get("仕様", "")
        target = base_info.get("ターゲット", "")
        scene = base_info.get("使用シーン", "")
        objective_a = base_info.get("目的", "")
        objective_b = base_info.get("戦略的目的", "")
        required_element = base_info.get("必須要素", "")
        concept_raw = json_data.get("採用コンセプト", "")
        tone_manner_raw = json_data.get("トーン_and_マナー規定", "")
    else:
        product_name = json_data.get("productName", "")
        item_name = json_data.get("itemName", "")
        spec = json_data.get("spec", "")
        target = json_data.get("target", "")
        scene = json_data.get("scene", "")
        objective_a = json_data.get("objectiveA", "")
        objective_b = json_data.get("objectiveB", "")
        required_element = json_data.get("requiredElement", "")
        concept_raw = json_data.get("concept", "")
        tm_data = json_data.get("toneManner", [])
        tone_manner_raw = "\n".join(tm_data) if isinstance(tm_data, list) else tm_data

    # ==========================================
    # ▼ 文章を綺麗に切り分ける処理（パース） ▼
    # ==========================================
    
    # 1. コンセプトの分割
    concept_title = concept_raw
    concept_desc = ""
    concept_design = ""
    
    # 「デザイン示唆：」で切り離す
    if "デザイン示唆：" in concept_title:
        parts = concept_title.split("デザイン示唆：", 1)
        concept_design = parts[1].strip()
        concept_title = parts[0]
        
    # 「説明：」で切り離す
    if "説明：" in concept_title:
        parts = concept_title.split("説明：", 1)
        concept_desc = parts[1].strip()
        concept_title = parts[0]
        
    # コンセプト名の不要文字を消す
    concept_title = concept_title.replace("コンセプト名：", "").replace("【案A】", "").replace("【案B】", "").replace("【案C】", "").replace("【案D】", "").replace("【案E】", "").strip()

    # 2. トーン＆マナーの分割
    tone_color = ""
    tone_photo = ""
    
    # 「写真イラスト：」で切り離す
    if "写真イラスト：" in tone_manner_raw:
        parts = tone_manner_raw.split("写真イラスト：", 1)
        tone_photo = parts[1].strip()
        tone_color = parts[0].replace("カラー：", "").strip()
    elif "写真イラスト" in tone_manner_raw: # コロンなし表記ゆれ対策
        parts = tone_manner_raw.split("写真イラスト", 1)
        tone_photo = parts[1].lstrip("：: ").strip()
        tone_color = parts[0].replace("カラー：", "").strip()
    else:
        tone_color = tone_manner_raw.replace("カラー：", "").strip()

    # ==========================================
    # ▲ 文章の切り分け処理ここまで ▲
    # ==========================================

    # パワポに書き込むための置換辞書
    replacements = {
        "{{productName}}": product_name,
        "{{itemName}}": item_name,
        "{{spec}}": spec,
        "{{target}}": target,
        "{{scene}}": scene,
        "{{objectiveA}}": objective_a,
        "{{objectiveB}}": objective_b,
        "{{requiredElement}}": required_element,
        # ▼ 新しく切り分けたタグ
        "{{concept}}": concept_title,
        "{{conceptDescription}}": concept_desc,
        "{{conceptDesign}}": concept_design,
        "{{toneColor}}": tone_color,
        "{{tonePhoto}}": tone_photo,
        # 古いテンプレート用の予備
        "{{toneManner}}": tone_manner_raw,
    }

    for slide in prs.slides:
        def process_shapes(shapes):
            for shape in shapes:
                if shape.shape_type == 6:
                    process_shapes(shape.shapes)
                elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                    replace_text_in_shape(shape, replacements)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            replace_text_in_shape(cell, replacements)
        process_shapes(slide.shapes)

    # 画像のスライド位置
    slide_indices = {"A案": 5, "B案": 6, "C案": 7}
    
    margin_x, margin_y = Inches(0.5), Inches(1.5)
    cell_w, cell_h = Inches(3.0), Inches(2.0)
    cols = 3

    for plan_name, images in uploaded_images.items():
        if plan_name in slide_indices and len(prs.slides) > slide_indices[plan_name]:
            slide = prs.slides[slide_indices[plan_name]]
            
            for idx, img_file in enumerate(images[:6]):
                row = idx // cols
                col = idx % cols
                x = margin_x + (col * cell_w)
                y = margin_y + (row * cell_h)
                
                img_stream = io.BytesIO(img_file.read())
                try:
                    slide.shapes.add_picture(img_stream, x, y, width=cell_w - Inches(0.2))
                except Exception as e:
                    st.warning(f"{plan_name}の画像挿入に失敗しました: {e}")

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# --- UI構築（左右2カラムレイアウト） ---
st.markdown("""
    <style>
        .block-container { padding-top: 1rem; padding-bottom: 1rem; }
        h1 { font-size: 1.6rem !important; margin-bottom: 1rem !important; }
        h2 { font-size: 1.2rem !important; margin-bottom: 0.2rem !important;}
        .stMarkdown p { font-size: 0.85rem; margin-bottom: 0.2rem !important;}
        [data-testid="stFileUploader"] { margin-bottom: 0rem; }
        [data-testid="stFileUploadDropzone"] {
            padding: 0.5rem !important;
            min-height: 1.5rem !important;
        }
        [data-testid="stFileUploadDropzone"] * {
            font-size: 0.8rem !important;
        }
        [data-testid="stFileUploadDropzone"] svg {
            display: none;
        }
        [data-testid="stFileUploader"] > section {
            max-height: 90px !important;
            overflow-y: auto !important;
        }
        [data-testid="stUploadedFile"] small {
            display: none !important;
        }
    </style>
""", unsafe_allow_html=True)

col1, col2 = st.columns(2, gap="large")

with col1:
    st.header("🖼️ 画像アップロード")
    st.markdown("各案の画像を枠内にドラッグ＆ドロップしてください。")

    uploaded_images = {}
    plans = ["A案", "B案", "C案"]

    for plan in plans:
        uploaded_images[plan] = st.file_uploader(
            f"📁 {plan}", 
            accept_multiple_files=True, 
            type=["png", "jpg", "jpeg"], 
            key=plan
        )

with col2:
    st.header("📝 企画書生成")
    st.markdown("左側のアプリからコピーしたJSONデータを貼り付けます。")

    json_text = st.text_area("JSONデータを貼り付け", height=280, label_visibility="collapsed", placeholder="ここにJSONデータを貼り付けてください")

    if st.button("📊 企画書パワーポイントを作成", type="primary", use_container_width=True):
        if not json_text.strip():
            st.error("エラー: JSONデータを入力してください。")
        else:
            try:
                json_data = json.loads(json_text)
                with st.spinner("PowerPointを生成中..."):
                    ppt_stream = generate_pptx(json_data, uploaded_images)
                    
                st.success("🎉 PowerPointの生成が完了しました！")
                item_name_from_json = json_data.get("基本情報", {}).get("アイテム名", "untitled")
                st.download_button(
                    label="📥 企画書(.pptx) をダウンロード",
                    data=ppt_stream,
                    file_name=f"proposal_{item_name_from_json}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
            except json.JSONDecodeError:
                st.error("エラー: JSONのフォーマットが正しくありません。")
            except Exception as e:
                st.error(f"予期せぬエラーが発生しました: {e}")
